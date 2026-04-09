"""
services.py — Business logic layer.

The main entry point is ``generate_documentation_stream()``, a generator
that yields SSE-ready event dicts as work progresses.  Routes consume it
via Flask's ``stream_with_context``.
"""

from __future__ import annotations

import os
import tkinter as tk
from pathlib import Path
from tkinter import filedialog

from . import md_to_docx
from .generators import DEFAULT_DOC_TYPE, get_generator
from .repo_reader import scan


# ── Backend message strings ───────────────────────────────────────────

_MSG = {
    "es": {
        "doc_type":          "Tipo de documento",
        "invalid_doc_type":  "Tipo de documento inválido",
        "scanning_repo":     "Analizando repositorio",
        "repo_read_error":   "No se pudo leer el repositorio",
        "repo_scan_error":   "Error inesperado al escanear el repositorio",
        "no_files":          (
            "No se encontraron archivos de código fuente en el repositorio. "
            "Verifica que la carpeta seleccionada sea la raíz del proyecto y "
            "que contenga archivos con extensiones reconocidas (.py, .js, .ts, etc.)."
        ),
        "files_found":       "Archivos encontrados",
        "files_skipped":     "archivo(s) omitido(s) (sin permisos de lectura o demasiado grandes).",
        "total":             "Total",
        "file_s":            ("archivo", "archivos"),
        "loading_template":  "Cargando plantilla",
        "template_ok":       "Plantilla cargada correctamente.",
        "building_prompt":   "Construyendo prompt...",
        "generating":        "Generando documentación. Esto puede tardar unos segundos...",
        "config_error":      "Error de configuración",
        "gen_error":         "Error al generar la documentación",
        "gen_unexpected":    "Error inesperado al llamar al modelo",
        "empty_response":    (
            "El modelo devolvió una respuesta vacía. "
            "Intenta de nuevo o revisa que el repositorio tenga contenido legible."
        ),
        "output_path_error": "No se pudo determinar la ruta de salida",
        "converting":        "Contenido generado. Convirtiendo a",
        "timeout_error":     (
            "La conversión superó el límite de {t} segundos. "
            "Intenta de nuevo o verifica que el repositorio no sea demasiado grande."
        ),
        "doc_error":         "Error al crear el documento",
        "export_error":      "Error inesperado al exportar el documento",
        "doc_ready":         "Documento listo",
        "no_repo":           "No se especificó ningún repositorio.",
    },
    "en": {
        "doc_type":          "Document type",
        "invalid_doc_type":  "Invalid document type",
        "scanning_repo":     "Scanning repository",
        "repo_read_error":   "Could not read the repository",
        "repo_scan_error":   "Unexpected error while scanning the repository",
        "no_files":          (
            "No source code files found in the repository. "
            "Make sure the selected folder is the project root and "
            "contains files with recognized extensions (.py, .js, .ts, etc.)."
        ),
        "files_found":       "Files found",
        "files_skipped":     "file(s) skipped (no read permission or too large).",
        "total":             "Total",
        "file_s":            ("file", "files"),
        "loading_template":  "Loading template",
        "template_ok":       "Template loaded successfully.",
        "building_prompt":   "Building prompt...",
        "generating":        "Generating documentation. This may take a few seconds...",
        "config_error":      "Configuration error",
        "gen_error":         "Error generating documentation",
        "gen_unexpected":    "Unexpected error calling the model",
        "empty_response":    (
            "The model returned an empty response. "
            "Try again or check that the repository has readable content."
        ),
        "output_path_error": "Could not determine the output path",
        "converting":        "Content generated. Converting to",
        "timeout_error":     (
            "Conversion exceeded the {t}-second limit. "
            "Try again or check that the repository is not too large."
        ),
        "doc_error":         "Error creating the document",
        "export_error":      "Unexpected error exporting the document",
        "doc_ready":         "Document ready",
        "no_repo":           "No repository was specified.",
    },
}


def _m(lang: str, key: str) -> str:
    """Return the message string for *key* in *lang*, falling back to Spanish."""
    return _MSG.get(lang, _MSG["es"]).get(key, _MSG["es"][key])


# ── SSE event constructors ────────────────────────────────────────────

def _log(message: str, level: str = "info") -> dict:
    return {"type": "log", "message": message, "level": level}


def _progress(pct: int) -> dict:
    return {"type": "progress", "pct": pct}


def _done(markdown: str, output_path: str) -> dict:
    return {"type": "done", "markdown": markdown, "output_path": output_path}


def _ready(output_path: str, filename: str) -> dict:
    return {"type": "ready", "output_path": output_path, "filename": filename}


def _error(message: str) -> dict:
    return {"type": "error", "message": message}


# ── Dialog helpers ────────────────────────────────────────────────────

def _get_tk_root():
    root = tk.Tk()
    root.withdraw()
    root.attributes("-topmost", True)
    return root


def browse_folder() -> str | None:
    root = _get_tk_root()
    path = filedialog.askdirectory(title="Seleccionar repositorio")
    root.destroy()
    return path or None


def browse_file() -> str | None:
    root = _get_tk_root()
    path = filedialog.askopenfilename(
        title="Seleccionar plantilla",
        filetypes=[
            ("Todos los archivos", "*.*"),
            ("Texto", "*.txt"),
            ("Markdown", "*.md"),
            ("Word", "*.docx"),
            ("PDF", "*.pdf"),
            ("PowerPoint", "*.pptx"),
        ],
    )
    root.destroy()
    return path or None


# ── Template section extractor ────────────────────────────────────────

def extract_template_sections(template_path: str) -> tuple[list[str], str | None]:
    """
    Read *template_path* and return a list of detected section headings.

    Returns (sections, error_message).  On success error_message is None.
    Sections are returned in document order; duplicates are removed while
    preserving order.

    Supported formats
    -----------------
    .md / .txt  — Markdown headings (lines starting with one or more '#')
    .docx       — Paragraphs whose Word style name begins with "Heading"
    .pdf        — Lines that match a heading pattern extracted via pypdf
    .pptx       — Slide title placeholder text extracted via python-pptx
    """
    p = Path(template_path)

    if not p.exists() or not p.is_file():
        return [], f"No se encontró el archivo: {template_path}"

    ext = p.suffix.lower()

    if ext == ".docx":
        try:
            from docx import Document as _Document
            wdoc = _Document(str(p))
            seen: list[str] = []
            for para in wdoc.paragraphs:
                style_name = (para.style.name or "").lower()
                text = para.text.strip()
                if style_name.startswith("heading") and text and text not in seen:
                    seen.append(text)
            return seen, None
        except Exception as exc:
            return [], f"No se pudo leer el archivo .docx: {exc}"

    if ext == ".pptx":
        try:
            from pptx import Presentation as _Prs
            from pptx.enum.shapes import PP_PLACEHOLDER
            prs  = _Prs(str(p))
            seen_set: set[str] = set()
            sections: list[str] = []
            for slide in prs.slides:
                for shape in slide.placeholders:
                    ph_type = shape.placeholder_format.type
                    if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                        title = shape.text.strip()
                        if title and title not in seen_set:
                            seen_set.add(title)
                            sections.append(title)
                        break
            return sections, None
        except ImportError:
            return [], "python-pptx no está instalado. Ejecuta: pip install python-pptx"
        except Exception as exc:
            return [], f"No se pudo leer el archivo .pptx: {exc}"

    if ext == ".pdf":
        try:
            from pypdf import PdfReader
            reader   = PdfReader(str(p))
            raw_text = "\n".join(
                (page.extract_text() or "") for page in reader.pages
            )
        except ImportError:
            return [], "pypdf no está instalado. Ejecuta: pip install pypdf"
        except Exception as exc:
            return [], f"No se pudo leer el archivo .pdf: {exc}"

        import re as _re
        heading_re = _re.compile(r"^#{1,4}\s+(.+)", _re.MULTILINE)
        seen_set_pdf: set[str] = set()
        sections_pdf: list[str] = []
        for m in heading_re.finditer(raw_text):
            title = m.group(1).strip()
            if title and title not in seen_set_pdf:
                seen_set_pdf.add(title)
                sections_pdf.append(title)

        # If no Markdown headings found, fall back to ALL-CAPS short lines as headings
        if not sections_pdf:
            for line in raw_text.splitlines():
                stripped = line.strip()
                if stripped and stripped.isupper() and 3 <= len(stripped) <= 80:
                    if stripped not in seen_set_pdf:
                        seen_set_pdf.add(stripped)
                        sections_pdf.append(stripped)

        return sections_pdf, None

    # .md / .txt — detect Markdown headings
    try:
        text = p.read_text(encoding="utf-8", errors="replace")
    except OSError as exc:
        return [], f"No se pudo leer el archivo: {exc}"

    import re
    heading_re = re.compile(r"^#{1,4}\s+(.+)", re.MULTILINE)
    seen_set_md: set[str] = set()
    sections_md: list[str] = []
    for m in heading_re.finditer(text):
        title = m.group(1).strip()
        if title and title not in seen_set_md:
            seen_set_md.add(title)
            sections_md.append(title)

    return sections_md, None


# ── Template reader ───────────────────────────────────────────────────

def _read_template(template_path: str) -> tuple[str | None, str | None]:
    """
    Read the template file and return (content, error_message).
    content is None when the file cannot be used; error_message explains why.
    """
    p = Path(template_path)

    if not p.exists():
        return None, f"El archivo de plantilla ya no existe: {template_path}"

    if not p.is_file():
        return None, f"La ruta de plantilla no apunta a un archivo: {template_path}"

    ext = p.suffix.lower()
    if ext == ".docx":
        return None, (
            "Las plantillas .docx no se pueden leer como texto plano todavía. "
            "Se usará la estructura predefinida."
        )

    try:
        content = p.read_text(encoding="utf-8", errors="strict")
        return content, None
    except UnicodeDecodeError:
        return None, (
            "La plantilla contiene caracteres que no se pueden leer como UTF-8. "
            "Se usará la estructura predefinida."
        )
    except PermissionError:
        return None, (
            f"Sin permisos para leer el archivo de plantilla: {template_path}. "
            "Se usará la estructura predefinida."
        )
    except OSError as exc:
        return None, (
            f"No se pudo leer el archivo de plantilla ({exc}). "
            "Se usará la estructura predefinida."
        )


# ── Main streaming generator ──────────────────────────────────────────

def generate_documentation_stream(
    repo_path: str,
    template_path: str | None = None,
    doc_type: str = DEFAULT_DOC_TYPE,
    primary_color: str | None = None,
    secondary_color: str | None = None,
    locked_sections: list[str] | None = None,
    section_enrichments: dict | None = None,
    api_key_override: str | None = None,
    model_override: str | None = None,
    provider_override: str | None = None,
    lang: str = "es",
):
    """
    Generator — yields SSE event dicts as work progresses.

    Sequence of events:
        log       — a line to display in the UI log area
        progress  — update the progress bar (0-100)
        done      — generation complete; carries markdown + output_path
        error     — fatal error; generation stops after this event
    """
    # Top-level guard: catch any bug we didn't anticipate so the stream
    # always closes cleanly instead of leaving the browser hanging.
    try:
        yield from _run(repo_path, template_path, doc_type, primary_color,
                        secondary_color, locked_sections, section_enrichments,
                        api_key_override, model_override, provider_override,
                        lang)
    except Exception as exc:
        yield _error(f"{'Unexpected internal error' if lang == 'en' else 'Error interno inesperado'}: {exc}.")


def _run(repo_path: str, template_path: str | None, doc_type: str,
         primary_color: str | None, secondary_color: str | None,
         locked_sections: list[str] | None = None,
         section_enrichments: dict | None = None,
         api_key_override: str | None = None,
         model_override: str | None = None,
         provider_override: str | None = None,
         lang: str = "es"):
    """Inner generator — all expected errors are handled here."""

    # ── 1. Validate inputs ───────────────────────────────────────────
    if not repo_path or not repo_path.strip():
        yield _error(_m(lang, "no_repo"))
        return

    try:
        generator = get_generator(doc_type)
    except ValueError as exc:
        yield _error(f"{_m(lang, 'invalid_doc_type')}: {exc}")
        return

    yield _log(f"{_m(lang, 'doc_type')}: {generator.DISPLAY_NAME}")
    yield _progress(5)

    # ── 2. Scan repository ───────────────────────────────────────────
    yield _log(f"{_m(lang, 'scanning_repo')}: {repo_path}")

    try:
        repo_scan = scan(repo_path)
    except ValueError as exc:
        yield _error(f"{_m(lang, 'repo_read_error')}. {exc}")
        return
    except Exception as exc:
        yield _error(f"{_m(lang, 'repo_scan_error')}: {exc}")
        return

    if repo_scan.total_files == 0:
        yield _error(_m(lang, "no_files"))
        return

    yield _log(f"{_m(lang, 'files_found')} ({repo_scan.total_files}):")
    for f in repo_scan.files:
        yield _log(f"  \u2192 {f.relative_path}  ({f.extension})")

    if repo_scan.skipped:
        yield _log(
            f"  {len(repo_scan.skipped)} {_m(lang, 'files_skipped')}",
            level="warn",
        )

    size_kb = repo_scan.total_bytes / 1024
    _fs = _m(lang, "file_s")
    _word = _fs[1] if repo_scan.total_files != 1 else _fs[0]
    yield _log(f"{_m(lang, 'total')}: {repo_scan.total_files} {_word} | {size_kb:.1f} KB")
    yield _progress(20)

    # ── 3. Load template ─────────────────────────────────────────────
    template_content: str | None = None
    if template_path:
        yield _log(f"{_m(lang, 'loading_template')}: {template_path}")
        template_content, err_msg = _read_template(template_path)
        if err_msg:
            yield _log(err_msg, level="warn")
        else:
            yield _log(_m(lang, "template_ok"), level="success")

    yield _progress(30)

    # ── 4. Generate Markdown via LLM ─────────────────────────────────
    try:
        repo_name = Path(repo_path).resolve().name
    except OSError:
        repo_name = Path(repo_path).name

    active_provider = (provider_override or "").strip() or "google"
    active_model    = (model_override or "").strip() or os.getenv("LLM_MODEL", "gemini-3-flash-preview").strip()
    yield _log(f"LLM: {active_provider} / {active_model}")
    yield _progress(40)

    yield _log(_m(lang, "building_prompt"))
    yield _log(_m(lang, "generating"))

    try:
        markdown = generator.generate(
            repo_scan, template_content, locked_sections,
            section_enrichments=section_enrichments,
            api_key_override=api_key_override,
            model_override=model_override,
            provider_override=provider_override,
        )
    except ValueError as exc:
        yield _error(f"{_m(lang, 'config_error')}: {exc}")
        return
    except RuntimeError as exc:
        yield _error(f"{_m(lang, 'gen_error')}: {exc}")
        return
    except Exception as exc:
        yield _error(f"{_m(lang, 'gen_unexpected')}: {exc}")
        return

    if not markdown or not markdown.strip():
        yield _error(_m(lang, "empty_response"))
        return

    yield _progress(90)

    # ── 5. Resolve output format + path ─────────────────────────────
    output_fmt = "docx"
    if template_path:
        _tpl_ext = Path(template_path).suffix.lower()
        if _tpl_ext == ".pdf":
            output_fmt = "pdf"
        elif _tpl_ext == ".pptx":
            output_fmt = "pptx"

    output_dir = (os.getenv("OUTPUT_DIR") or "").strip() or repo_path

    try:
        output_path = str(generator.output_path(repo_name, output_dir, fmt=output_fmt))
    except Exception as exc:
        yield _error(f"{_m(lang, 'output_path_error')}: {exc}")
        return

    _fmt_labels = {"docx": "Word (.docx)", "pdf": "PDF (.pdf)", "pptx": "PowerPoint (.pptx)"}
    yield _log(
        f"{_m(lang, 'converting')} {_fmt_labels.get(output_fmt, output_fmt)}...",
        level="success",
    )
    yield _progress(93)

    # ── 6. Convert Markdown → output format ──────────────────────────
    import concurrent.futures as _cf

    _CONVERT_TIMEOUT = 60  # seconds

    kwargs: dict = {"doc_type": doc_type or "technical"}
    if primary_color:
        kwargs["primary_color"] = primary_color
    if secondary_color:
        kwargs["secondary_color"] = secondary_color

    if output_fmt == "pptx":
        from . import md_to_pptx
        _pptx_kwargs = {k: v for k, v in kwargs.items() if k in ("primary_color", "secondary_color")}
        _converter = lambda: md_to_pptx.convert(markdown, output_path, **_pptx_kwargs)  # noqa: E731
    elif output_fmt == "pdf":
        from . import md_to_pdf
        _converter = lambda: md_to_pdf.convert(markdown, output_path, **kwargs)  # noqa: E731
    else:
        _converter = lambda: md_to_docx.convert(markdown, output_path, **kwargs)  # noqa: E731

    try:
        with _cf.ThreadPoolExecutor(max_workers=1) as _pool:
            _future = _pool.submit(_converter)
            try:
                final_path = _future.result(timeout=_CONVERT_TIMEOUT)
            except _cf.TimeoutError:
                _future.cancel()
                yield _error(_m(lang, "timeout_error").format(t=_CONVERT_TIMEOUT))
                return
    except RuntimeError as exc:
        yield _error(f"{_m(lang, 'doc_error')}: {exc}")
        return
    except Exception as exc:
        yield _error(f"{_m(lang, 'export_error')}: {exc}")
        return

    filename = Path(output_path).name
    yield _log(f"{_m(lang, 'doc_ready')}: {filename}", level="success")
    yield _progress(100)
    yield _ready(output_path=str(final_path), filename=filename)
