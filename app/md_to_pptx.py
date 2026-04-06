"""
md_to_pptx.py — Convert a Markdown string to a PowerPoint (.pptx) presentation.

Each ``##`` heading becomes a content slide.
The first ``#`` heading becomes the cover slide title.
Sub-headings (###, ####) are treated as bold paragraph text within the slide.

Public API
----------
convert(markdown, output_path, *, primary_color, secondary_color) -> Path
"""

from __future__ import annotations

import re
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── Defaults ──────────────────────────────────────────────────────────

_DEFAULT_PRIMARY   = "#2D3A8C"
_DEFAULT_SECONDARY = "#287A5F"

_MONTHS_ES = [
    "enero", "febrero", "marzo", "abril", "mayo", "junio",
    "julio", "agosto", "septiembre", "octubre", "noviembre", "diciembre",
]


def _month_year_es() -> str:
    now = datetime.now()
    return f"{_MONTHS_ES[now.month - 1].capitalize()} {now.year}"


def _hex_to_pptx(hex_color: str) -> PptxRGB:
    h = hex_color.lstrip("#")
    return PptxRGB(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ── Markdown parser ───────────────────────────────────────────────────

def _parse(markdown: str) -> tuple[str, list[tuple[str, list[str]]]]:
    """
    Parse *markdown* into (title, [(section_heading, content_lines)]).

    - The first ``# heading`` becomes *title*.
    - Each ``## heading`` (and deeper) opens a new section.
    - Lines inside fenced code blocks and [DIAGRAM] blocks are skipped.
    """
    lines      = markdown.splitlines()
    title      = ""
    sections: list[tuple[str, list[str]]] = []
    cur_head: str | None = None
    cur_lines: list[str] = []
    in_code    = False
    in_diagram = False

    for line in lines:
        stripped = line.strip()

        # Track fenced code / diagram blocks
        if stripped.startswith("```"):
            in_code = not in_code
            continue
        if stripped == "[DIAGRAM]":
            in_diagram = True
            continue
        if stripped == "[/DIAGRAM]":
            in_diagram = False
            continue
        if in_code or in_diagram:
            continue

        # Skip caption tags
        if stripped.startswith("[CAPTION:") and stripped.endswith("]"):
            continue

        # H1 → cover title (only the first one)
        h1 = re.match(r"^#\s+(.*)", line)
        if h1 and not title:
            title = h1.group(1).strip()
            continue

        # H2+ → new section
        h2 = re.match(r"^#{2,}\s+(.*)", line)
        if h2:
            if cur_head is not None:
                sections.append((cur_head, cur_lines))
            cur_head  = h2.group(1).strip()
            cur_lines = []
            continue

        if cur_head is not None and stripped:
            cur_lines.append(line)

    if cur_head is not None:
        sections.append((cur_head, cur_lines))

    if not title:
        title = "Documentación"

    return title, sections


# ── Slide builders ────────────────────────────────────────────────────

# Widescreen 16:9 (default in python-pptx)
_W = Inches(13.333)
_H = Inches(7.5)

_HEADER_H  = Inches(1.25)
_MARGIN    = Inches(0.45)
_BODY_TOP  = _HEADER_H + Inches(0.25)
_BODY_H    = _H - _BODY_TOP - _MARGIN


def _add_cover_slide(prs: Presentation, title: str, subtitle: str,
                     primary: PptxRGB, secondary: PptxRGB) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Background rectangle (primary color, top half)
    bg = slide.shapes.add_shape(1, 0, 0, _W, Inches(3.8))  # MSO_SHAPE_TYPE.RECTANGLE = 1
    bg.fill.solid()
    bg.fill.fore_color.rgb = primary
    bg.line.fill.background()

    # Title text
    tf_title = slide.shapes.add_textbox(
        _MARGIN, Inches(0.9), _W - _MARGIN * 2, Inches(2.2)
    ).text_frame
    tf_title.word_wrap = True
    p = tf_title.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text          = title
    run.font.bold     = True
    run.font.size     = Pt(40)
    run.font.color.rgb = PptxRGB(0xFF, 0xFF, 0xFF)

    # Subtitle / date block (below the colored area)
    tf_sub = slide.shapes.add_textbox(
        _MARGIN, Inches(4.1), _W - _MARGIN * 2, Inches(1.2)
    ).text_frame
    p2 = tf_sub.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text           = subtitle
    r2.font.size      = Pt(18)
    r2.font.color.rgb = secondary

    # Date bottom-right
    tf_date = slide.shapes.add_textbox(
        _W - Inches(3.5), _H - Inches(0.65), Inches(3.2), Inches(0.45)
    ).text_frame
    pd = tf_date.paragraphs[0]
    pd.alignment = PP_ALIGN.RIGHT
    rd = pd.add_run()
    rd.text           = _month_year_es()
    rd.font.size      = Pt(11)
    rd.font.color.rgb = PptxRGB(0x88, 0x88, 0x99)


def _add_content_slide(prs: Presentation, heading: str, content_lines: list[str],
                       primary: PptxRGB, secondary: PptxRGB) -> None:
    slide = prs.shapes  # not used; we need prs.slides
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Header bar
    hdr = slide.shapes.add_shape(1, 0, 0, _W, _HEADER_H)
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = primary
    hdr.line.fill.background()

    # Heading text in header
    tf_hdr = slide.shapes.add_textbox(
        _MARGIN, Inches(0.25), _W - _MARGIN * 2, Inches(0.85)
    ).text_frame
    ph = tf_hdr.paragraphs[0]
    ph.alignment = PP_ALIGN.LEFT
    rh = ph.add_run()
    rh.text           = heading
    rh.font.bold      = True
    rh.font.size      = Pt(24)
    rh.font.color.rgb = PptxRGB(0xFF, 0xFF, 0xFF)

    # Content area
    tf_body = slide.shapes.add_textbox(
        _MARGIN, _BODY_TOP, _W - _MARGIN * 2, _BODY_H
    ).text_frame
    tf_body.word_wrap = True

    first_para = True
    for line in content_lines:
        stripped = line.strip()
        if not stripped:
            continue

        para = tf_body.paragraphs[0] if first_para else tf_body.add_paragraph()
        first_para = False

        # Bullet list item
        bullet_m = re.match(r"^\s*[-*]\s+(.*)", line)
        num_m    = re.match(r"^\s*\d+\.\s+(.*)", line)
        # Sub-heading (###) → bold paragraph
        sub_m    = re.match(r"^#{3,}\s+(.*)", line)

        if sub_m:
            para.alignment = PP_ALIGN.LEFT
            run = para.add_run()
            run.text           = sub_m.group(1).strip()
            run.font.bold      = True
            run.font.size      = Pt(14)
            run.font.color.rgb = secondary
        elif bullet_m or num_m:
            text = (bullet_m or num_m).group(1).strip()
            para.alignment = PP_ALIGN.LEFT
            # indent via space prefix (python-pptx paragraph level)
            para.level = 1
            run = para.add_run()
            run.text      = text
            run.font.size = Pt(13)
            run.font.color.rgb = PptxRGB(0x1D, 0x1D, 0x1F)
        else:
            para.alignment = PP_ALIGN.LEFT
            # Parse inline **bold** simply
            parts = re.split(r"\*\*(.+?)\*\*", stripped)
            bold_next = False
            for i, part in enumerate(parts):
                if not part:
                    continue
                run = para.add_run()
                run.text      = part
                run.font.size = Pt(13)
                run.font.color.rgb = PptxRGB(0x1D, 0x1D, 0x1F)
                run.font.bold = (i % 2 == 1)  # odd chunks are inside **...**


# ── Public API ────────────────────────────────────────────────────────

def convert(
    markdown: str,
    output_path: str | Path,
    *,
    primary_color:   str = _DEFAULT_PRIMARY,
    secondary_color: str = _DEFAULT_SECONDARY,
) -> Path:
    """
    Convert *markdown* to a .pptx file at *output_path*.

    Returns the resolved Path of the written file.

    Raises
    ------
    RuntimeError — if the output directory does not exist.
    """
    out = Path(output_path)
    if not out.parent.exists():
        raise RuntimeError(
            f"El directorio de salida no existe: {out.parent}."
        )

    primary   = _hex_to_pptx(primary_color)
    secondary = _hex_to_pptx(secondary_color)

    title, sections = _parse(markdown)

    prs = Presentation()
    prs.slide_width  = _W
    prs.slide_height = _H

    # Cover slide
    doc_subtitle = "Draft generado por GenDoc"
    _add_cover_slide(prs, title, doc_subtitle, primary, secondary)

    # One content slide per section
    for heading, content_lines in sections:
        _add_content_slide(prs, heading, content_lines, primary, secondary)

    prs.save(str(out))
    return out.resolve()
